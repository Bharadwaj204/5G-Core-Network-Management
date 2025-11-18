"""
Python script to generate PowerPoint presentation for 5G Core Management Prototype
Uses python-pptx library to create a professional presentation with images
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from typing import Optional

def create_presentation():
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define styles
    def set_title_style(title):
        if title and hasattr(title, 'text_frame') and title.text_frame:
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        
    def set_content_style(content):
        if content and hasattr(content, 'text_frame') and content.text_frame:
            for paragraph in content.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(0, 0, 0)
    
    # Function to add image to slide if it exists
    def add_image_if_exists(slide, image_path, left, top, width, height):
        if os.path.exists(image_path):
            try:
                slide.shapes.add_picture(image_path, left, top, width, height)
                return True
            except Exception as e:
                print(f"Warning: Could not add image {image_path}: {e}")
                return False
        else:
            print(f"Warning: Image not found: {image_path}")
            return False
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    try:
        if title:
            title.text = "End-to-End 5G Core Management Prototype"
        set_title_style(title)
        
        if subtitle:
            subtitle.text = "A Hands-on 5G Core Management System\nUsing NETCONF, RESTCONF, SNMP, YANG, and Open5GS"
            if hasattr(subtitle, 'text_frame') and subtitle.text_frame:
                subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    except AttributeError:
        pass
    
    # Slide 2: Problem Statement
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    try:
        if title:
            title.text = "Problem Statement"
        set_title_style(title)
        
        if content:
            content.text = (
                "• Students have theoretical knowledge of 5G Core, but no hands-on platform to learn management\n"
                "• Real telecom systems are expensive and complex\n"
                "• Need a lightweight, open-source, manageable 5G Core"
            )
        set_content_style(content)
    except AttributeError:
        pass
    
    # Slide 3: Project Objective
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    try:
        if title:
            title.text = "Project Objective"
        set_title_style(title)
        
        if content:
            content.text = (
                "• Deploy Open5GS (a free 5G Core implementation)\n"
                "• Manage it using industry-standard protocols:\n"
                "  - NETCONF\n"
                "  - RESTCONF\n"
                "  - SNMP\n"
                "• Build YANG models\n"
                "• Create dashboards and monitoring\n"
                "• Provide easy test scripts"
            )
        set_content_style(content)
    except AttributeError:
        pass
    
    # Slide 4: 5G Core Architecture
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    try:
        if title:
            title.text = "5G Core Architecture"
        set_title_style(title)
        
        if content:
            content.text = (
                "• AMF (Access and Mobility Management Function)\n"
                "• SMF (Session Management Function)\n"
                "• UPF (User Plane Function)\n"
                "• NRF, UDM, AUSF\n"
                "• N2/N3 interfaces\n"
                "• UE/gNB (User Equipment and Base Station)"
            )
        set_content_style(content)
        
        # Try to add architecture diagram
        diagram_path = "architecture-diagram.png"
        if os.path.exists(diagram_path):
            add_image_if_exists(slide, diagram_path, Inches(7), Inches(2), Inches(5), Inches(3))
    except AttributeError:
        pass
    
    # Slide 5: System Architecture
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    try:
        if title:
            title.text = "What We Built (System Architecture)"
        set_title_style(title)
        
        if content:
            content.text = (
                "Two main layers:\n\n"
                "1. 5G Core Plane (Open5GS: AMF/SMF/UPF/etc.)\n\n"
                "2. Management Plane:\n"
                "   • NETCONF Server\n"
                "   • RESTCONF API\n"
                "   • SNMP Monitor\n"
                "   • YANG Models\n"
                "   • Dashboard"
            )
        set_content_style(content)
    except AttributeError:
        pass
    
    # Slide 6: Management Technologies
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    try:
        if title:
            title.text = "Management Technologies"
        set_title_style(title)
        
        if content:
            content.text = (
                "1. YANG Models\n"
                "   • Data modeling language\n"
                "   • Defines structure of subscribers, sessions, QoS profiles\n\n"
                "2. NETCONF\n"
                "   • XML, SSH-based network configuration\n"
                "   • Used by telecom vendors (Nokia, Ericsson, Huawei)\n\n"
                "3. RESTCONF\n"
                "   • HTTP/JSON version of NETCONF\n"
                "   • Lightweight, easier for developers\n\n"
                "4. SNMP\n"
                "   • Used for monitoring only\n"
                "   • Metrics and performance statistics"
            )
        set_content_style(content)
    except AttributeError:
        pass
    
    # Slide 7: YANG Models Created
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    try:
        if title:
            title.text = "YANG Models Created"
        set_title_style(title)
        
        if content:
            content.text = (
                "• subscribers.yang\n"
                "• sessions.yang\n"
                "• qos-profiles.yang\n"
                "• network-functions.yang\n\n"
                "What they represent:\n"
                "• subscribers.yang - Subscriber information (IMSI, MSISDN, security)\n"
                "• sessions.yang - PDU session details (QoS, DNN, status)\n"
                "• qos-profiles.yang - Quality of Service parameters\n"
                "• network-functions.yang - Core network functions (AMF, SMF, UPF)\n\n"
                "Example snippet:\n"
                "container subscribers {\n"
                "   list subscriber {\n"
                "      key \"imsi\"\n"
                "      leaf imsi { type string; }\n"
                "   }\n"
                "}"
            )
            if hasattr(content, 'text_frame') and content.text_frame:
                content.text_frame.paragraphs[-1].font.size = Pt(14)  # Smaller font for code
        set_content_style(content)
    except AttributeError:
        pass
    
    # Slide 8: NETCONF Implementation
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    try:
        if title:
            title.text = "NETCONF Implementation"
        set_title_style(title)
        
        if content:
            content.text = (
                "• Uses ncclient library\n"
                "• Runs on port 830\n"
                "• Supports get, edit-config, get-config\n"
                "• Stores data using YANG models\n\n"
                "Sample command:\n"
                "from ncclient import manager\n"
                "m = manager.connect(host='localhost', port=830, username='admin', password='admin')\n"
                "config = m.get_config(source=\"running\")\n"
                "print(config.data_xml)"
            )
            if hasattr(content, 'text_frame') and content.text_frame:
                content.text_frame.paragraphs[-1].font.size = Pt(14)  # Smaller font for code
        set_content_style(content)
    except AttributeError:
        pass
    
    # Slide 9: RESTCONF Implementation
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    try:
        if title:
            title.text = "RESTCONF Implementation"
        set_title_style(title)
        
        if content:
            content.text = (
                "• Built using Flask\n"
                "• Uses same YANG data structures\n"
                "• Key endpoints created:\n"
                "  /restconf/data/subscribers\n"
                "  /restconf/data/sessions\n"
                "  /restconf/data/qos-profiles\n"
                "  /restconf/data/network-functions\n\n"
                "RESTCONF = NETCONF + HTTP + JSON"
            )
        set_content_style(content)
    except AttributeError:
        pass
    
    # Slide 10: SNMP Monitoring
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    try:
        if title:
            title.text = "SNMP Monitoring"
        set_title_style(title)
        
        if content:
            content.text = (
                "• Uses PySNMP library\n"
                "• Collects metrics from Open5GS\n"
                "• Exports to JSON/CSV formats\n"
                "• Metrics include:\n"
                "  - Active sessions\n"
                "  - CPU utilization\n"
                "  - Memory usage\n"
                "  - Throughput statistics\n"
                "  - Packet loss rates"
            )
        set_content_style(content)
    except AttributeError:
        pass
    
    # Slide 11: Dashboard
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    try:
        if title:
            title.text = "Dashboard"
        set_title_style(title)
        
        if content:
            content.text = (
                "• Real-time metrics visualization\n"
                "• Subscriber/session count monitoring\n"
                "• Performance graphs\n"
                "• Built with Node.js and Socket.IO\n"
                "• Integration with Grafana for professional monitoring"
            )
        set_content_style(content)
    except AttributeError:
        pass
    
    # Slide 12: Project Folder Structure
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    try:
        if title:
            title.text = "Project Folder Structure"
        set_title_style(title)
        
        if content:
            content.text = (
                "5g-core-management/\n"
                "├── open5gs-setup/\n"
                "├── management-plane/\n"
                "│   ├── netconf-server/\n"
                "│   ├── restconf-api/\n"
                "│   ├── snmp-monitor/\n"
                "│   └── yang-models/\n"
                "├── dashboard/\n"
                "├── docs/\n"
                "├── tests/\n"
                "└── start_system.py"
            )
            if hasattr(content, 'text_frame') and content.text_frame:
                content.text_frame.paragraphs[0].font.size = Pt(14)  # Smaller font for tree
        set_content_style(content)
    except AttributeError:
        pass
    
    # Slide 13: Demo Flow
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    try:
        if title:
            title.text = "Demo Flow"
        set_title_style(title)
        
        if content:
            content.text = (
                "1. Starting NETCONF server\n"
                "   cd management-plane/netconf-server\n"
                "   python netconf_server.py\n\n"
                "2. NETCONF Get-Config Demo\n"
                "   python scripts/get_subscribers.py\n\n"
                "3. RESTCONF Demo\n"
                "   GET http://localhost:830/restconf/data/subscribers\n\n"
                "4. SNMP Demo\n"
                "   snmpwalk -v2c -c public localhost .1.3.6.1.4.1.55555\n\n"
                "5. Dashboard Demo\n"
                "   Show real-time updating graphs"
            )
            if hasattr(content, 'text_frame') and content.text_frame:
                content.text_frame.paragraphs[0].font.size = Pt(14)  # Smaller font for commands
        set_content_style(content)
    except AttributeError:
        pass
    
    # Slide 14: Results
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    try:
        if title:
            title.text = "Results"
        set_title_style(title)
        
        if content:
            content.text = (
                "✅ Open5GS deployed successfully\n"
                "✅ NETCONF + RESTCONF interfaces working\n"
                "✅ SNMP monitoring functioning\n"
                "✅ Dashboard showing real-time metrics\n"
                "✅ End-to-end 5G management achieved"
            )
        set_content_style(content)
    except AttributeError:
        pass
    
    # Slide 15: Challenges Faced
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    try:
        if title:
            title.text = "Challenges Faced"
        set_title_style(title)
        
        if content:
            content.text = (
                "• Open5GS configuration issues\n"
                "• YANG model validation complexities\n"
                "• NETCONF XML formatting requirements\n"
                "• SNMP OID management\n"
                "• Synchronizing NETCONF ↔ RESTCONF states\n"
                "• Python library compatibility issues"
            )
        set_content_style(content)
    except AttributeError:
        pass
    
    # Slide 16: Conclusion
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    try:
        if title:
            title.text = "Conclusion"
        set_title_style(title)
        
        if content:
            content.text = (
                "• Built a complete 5G Management Plane\n"
                "• Used real telecom protocols\n"
                "• Provides actual hands-on environment\n"
                "• Fully aligned with industry standards\n"
                "• Demonstrates practical management capabilities"
            )
        set_content_style(content)
    except AttributeError:
        pass
    
    # Slide 17: Future Enhancements
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    try:
        if title:
            title.text = "Future Enhancements"
        set_title_style(title)
        
        if content:
            content.text = (
                "• gNB + UE integration for complete end-to-end testing\n"
                "• Complete GUI-based provisioning system\n"
                "• Alarm management system\n"
                "• TR-069 integration for device management\n"
                "• Enhanced security (TLS, token authentication)\n"
                "• Containerization with Docker/Kubernetes"
            )
        set_content_style(content)
    except AttributeError:
        pass
    
    # Slide 18: Thank You
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    try:
        if title:
            title.text = "Thank You"
        set_title_style(title)
        
        if subtitle:
            subtitle.text = "Questions?"
            if hasattr(subtitle, 'text_frame') and subtitle.text_frame:
                subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    except AttributeError:
        pass
    
    # Save presentation
    prs.save("5G-Core-Management-Enhanced-Presentation.pptx")
    print("Enhanced presentation with images created successfully!")

if __name__ == "__main__":
    create_presentation()